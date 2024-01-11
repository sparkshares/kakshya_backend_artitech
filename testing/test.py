import os
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
from PIL import Image
import pytesseract
import cv2
import nltk
from collections import Counter
from nltk.tokenize import word_tokenize, sent_tokenize
from nltk.corpus import stopwords, words
from string import punctuation

# Ensure NLTK components are downloaded
nltk.download('punkt')
nltk.download('stopwords')
nltk.download('words')

def preprocess_image(image_path):
    img = cv2.imread(image_path, cv2.IMREAD_GRAYSCALE)
    img = cv2.adaptiveThreshold(img, 255, cv2.ADAPTIVE_THRESH_GAUSSIAN_C, cv2.THRESH_BINARY, 11, 2)
    preprocessed_image_path = 'preprocessed.png'
    cv2.imwrite(preprocessed_image_path, img)
    return preprocessed_image_path

def extract_text_from_file(file_path):
    if file_path.endswith('.pdf'):
        with open(file_path, 'rb') as file:
            pdf_reader = PdfReader(file)
            text = ''
            for page in pdf_reader.pages:
                extracted_text = page.extract_text()
                if extracted_text:
                    text += extracted_text
        return text
    elif file_path.endswith('.docx'):
        doc = Document(file_path)
        text = ' '.join([paragraph.text for paragraph in doc.paragraphs])
        return text
    elif file_path.endswith('.pptx'):
        pres = Presentation(file_path)
        text = ''
        for slide in pres.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            text += run.text
        return text
    elif file_path.endswith(('.png', '.jpg', '.jpeg')):
        preprocessed_image_path = preprocess_image(file_path)
        image = Image.open(preprocessed_image_path)
        config = '-c preserve_interword_spaces=1 --oem 3 --psm 6'
        text = pytesseract.image_to_string(image, config=config)
        return text
    else:
        return "Unsupported file type"

def insert_spaces(text):
    english_vocab = set(words.words())
    start = 0
    end = 2
    result = ""
    while start < len(text):
        if text[start:end] in english_vocab or end == len(text):
            result += text[start:end] + " "
            start = end
            end = start + 2
        else:
            end += 1
    return result

def summarize_text(text, num_sentences=10):
    words = word_tokenize(text.lower())
    stop_words = set(stopwords.words('english'))
    words = [word for word in words if word not in stop_words and word not in punctuation]
    freq = Counter(words)
    sentences = sent_tokenize(text)
    sentence_scores = {sentence: sum(freq.get(word, 0) for word in word_tokenize(sentence)) for sentence in sentences}
    summary_sentences = sorted(sentence_scores, key=sentence_scores.get, reverse=True)
    summary = ' '.join(summary_sentences[:num_sentences])
    return summary

def process_and_summarize_file(file_path):
    if not os.path.exists(file_path):
        print("File does not exist.")
        return

    extracted_text = extract_text_from_file(file_path)

    if not extracted_text:
        print("No text extracted from the file.")
        return

    segmented_text = insert_spaces(extracted_text)
    summarized_text = summarize_text(segmented_text)

    with open('summary_output.txt', 'w') as file:
        file.write(summarized_text)

# Example usage:
process_and_summarize_file('testing/chapter1.pdf')
