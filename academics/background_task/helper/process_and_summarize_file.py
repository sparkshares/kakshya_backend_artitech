import os
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
from io import BytesIO

def extract_text_from_file(file_path):
    """
    Extracts text from a given file.
    Supported file types: PDF, DOCX, PPTX.
    """
    text = ''
    if file_path.endswith('.pdf'):
        with open(file_path, 'rb') as file:
            pdf_reader = PdfReader(file)
            for page in pdf_reader.pages:
                extracted_text = page.extract_text()
                if extracted_text:
                    text += extracted_text + "\n"
                if len(text) > 50000:
                    return text[:50000]
    elif file_path.endswith('.docx'):
        doc = Document(file_path)
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
            if len(text) > 50000:
                return text[:50000]
    elif file_path.endswith('.pptx'):
        pres = Presentation(file_path)
        for slide in pres.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            text += run.text + "\n"
                if len(text) > 50000:
                    return text[:50000]
    else:
        print("Unsupported file type")
        return None
    return text

def process_file(file_path):
    """
    Processes the given file and returns its text.
    """
    if not os.path.exists(file_path):
        print("File does not exist.")
        return None

    return extract_text_from_file(file_path)